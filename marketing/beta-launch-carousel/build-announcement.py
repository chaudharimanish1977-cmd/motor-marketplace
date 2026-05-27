"""
The Beta-launch ANNOUNCEMENT deck — the 5-slide founder-voice post.

Distinct from the 10-slide deep deck (build.py). This is what we post
on LinkedIn day 1. The 10-slider becomes the website / follow-up
artefact.

Narrative arc:
  1. HOOK            — pattern-interrupt opener
  2. WHY NOW         — recognition + 'why we built this' in we-voice
  3. PROOF           — illustrative gap finding, clearly flagged
  4. ASK             — forward to review@rightoffer.in
  5. WHAT'S NEXT     — soft Phase 2 signal + wordmark

Same visual system as build.py (dark mode, 1080×1080) — re-uses every
helper + asset.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Import everything from the deep-deck build so we share helpers,
# colours, fonts, and the prerender step.
from build import (
    BG_DARK, TEXT_LIGHT, PLUM_LIFT, SAGE_LIFT, SLATE_LIFT, HAIRLINE,
    SERIF, MONO, CANVAS_W, CANVAS_H,
    add_bg, add_text, set_run, add_kicker, add_closer,
    add_wordmark, add_pic, add_paragraph,
    prerender_sketches, _asset_path,
)


# ─── Slides ───────────────────────────────────────────────────────────

def slide_hook(prs):
    """Pattern-interrupt opener.

    Headline reads as one sentence, but visually 'once.' is the
    punchline — italic plum, oversized — so the eye lands on it.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  beta is live  · · ·", Inches(1.0))

    # Line 1 — "Most owners read their policy"
    tb, tf = add_text(s, Inches(0.6), Inches(2.4), CANVAS_W - Inches(1.2),
                      Inches(1.4), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.1
    r = p.add_run()
    set_run(r, text="Most owners read their policy",
            font=SERIF, size=46, color=TEXT_LIGHT)

    # Line 2 — italic plum punchline (HUGE)
    tb, tf = add_text(s, Inches(0.6), Inches(3.8), CANVAS_W - Inches(1.2),
                      Inches(2.4), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="once.", font=SERIF, size=140, color=PLUM_LIFT,
            italic=True, bold=True)

    # Sage underline beneath 'once.' for spotlight
    underline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(4.05), Inches(6.4),
                                   Inches(2.7), Emu(48000))
    underline.fill.solid(); underline.fill.fore_color.rgb = SAGE_LIFT
    underline.line.fill.background(); underline.shadow.inherit = False

    # Subhead
    tb, tf = add_text(s, Inches(0.6), Inches(7.0), CANVAS_W - Inches(1.2),
                      Inches(0.9), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="After it’s too late.", font=SERIF, size=30,
            color=TEXT_LIGHT, italic=True)

    add_wordmark(s, center_y=Inches(9.9))


def slide_why_now(prs):
    """Why-now / why-this story in 'we' voice — family-burned angle.

    No diagrams — just a calm, declarative paragraph + a small
    supporting sketch. Reads like the founder is talking to you.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  why we built this  · · ·", Inches(1.0))

    # Story — first paragraph
    tb, tf = add_text(s, Inches(0.9), Inches(2.4), CANVAS_W - Inches(1.8),
                      Inches(3.4), align="center", anchor="middle")
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.35
    r = p.add_run()
    set_run(r, text="Someone close to us discovered ",
            font=SERIF, size=32, color=TEXT_LIGHT)
    r = p.add_run()
    set_run(r, text="at claim time", font=SERIF, size=32,
            color=PLUM_LIFT, italic=True)
    r = p.add_run()
    set_run(r, text=" that their motor policy didn’t actually cover their car.",
            font=SERIF, size=32, color=TEXT_LIGHT)

    # Beat
    p2 = add_paragraph(tf, alignment="center"); p2.line_spacing = 1.35
    r = p2.add_run()
    set_run(r, text=" ", font=SERIF, size=18, color=TEXT_LIGHT)

    p3 = add_paragraph(tf, alignment="center"); p3.line_spacing = 1.35
    r = p3.add_run()
    set_run(r, text="So we built the audit nobody else was going to.",
            font=SERIF, size=26, color=TEXT_LIGHT, italic=True)

    # Loupe sketch sits below — visual of 'reading carefully'
    add_pic(s, "loupe", left=Inches(4.05), top=Inches(7.2),
            width=Inches(2.7))


def slide_proof(prs):
    """A single illustrative gap finding rendered as a card.

    Clear visual hierarchy:
      KICKER  · · · GAP CAUGHT · · ·
      LABEL   Zero-depreciation cover
      VERB    expired silently.
      AMOUNT  ~ ₹ 38,000
      CAPTION What it would have cost on a routine accident claim.
      FOOTER  Representative audit finding — anonymised.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  gap caught  · · ·", Inches(1.0))

    # Finding label
    tb, tf = add_text(s, Inches(0.6), Inches(2.4), CANVAS_W - Inches(1.2),
                      Inches(0.8), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="Zero-depreciation cover", font=SERIF, size=34,
            color=TEXT_LIGHT)

    # Verb / state
    tb, tf = add_text(s, Inches(0.6), Inches(3.2), CANVAS_W - Inches(1.2),
                      Inches(1.0), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="expired silently.", font=SERIF, size=44,
            color=PLUM_LIFT, italic=True, bold=True)

    # Big amount — the punch
    tb, tf = add_text(s, Inches(0.5), Inches(4.6), CANVAS_W - Inches(1.0),
                      Inches(2.0), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="~ ₹38,000", font=SERIF, size=96, color=TEXT_LIGHT,
            bold=True)

    # Caption
    tb, tf = add_text(s, Inches(1.0), Inches(7.0), CANVAS_W - Inches(2.0),
                      Inches(1.0), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.4
    r = p.add_run()
    set_run(r,
            text="What it would have cost on a routine accident claim.",
            font=SERIF, size=20, color=TEXT_LIGHT, italic=True)

    # Hairline
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(4.0), Inches(8.5),
                              Inches(2.8), Emu(12700))
    line.fill.solid(); line.fill.fore_color.rgb = SAGE_LIFT
    line.line.fill.background(); line.shadow.inherit = False

    # Footer — clearly flag as illustrative
    tb, tf = add_text(s, Inches(0.6), Inches(8.8), CANVAS_W - Inches(1.2),
                      Inches(0.7), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r,
            text="Representative audit finding · anonymised composite",
            font=MONO, size=11, color=SLATE_LIFT, spacing=3)


def slide_ask(prs):
    """The action. Email primary, large, sage underline."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  forward your policy  · · ·", Inches(1.0))

    # Envelope hero
    add_pic(s, "envelope", left=Inches(4.0), top=Inches(2.1),
            width=Inches(2.8))

    # "FORWARD TO" tag
    tb, tf = add_text(s, Inches(0.5), Inches(4.9), CANVAS_W - Inches(1.0),
                      Inches(0.5), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="FORWARD TO", font=MONO, size=13, color=SAGE_LIFT,
            spacing=5, bold=True)

    # Email — big plum italic
    tb, tf = add_text(s, Inches(0.4), Inches(5.5), CANVAS_W - Inches(0.8),
                      Inches(1.5), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="review@rightoffer.in", font=SERIF, size=46,
            color=PLUM_LIFT, italic=True, bold=True)

    # Sage underline
    underline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1.5), Inches(6.65),
                                   Inches(7.8), Emu(36000))
    underline.fill.solid(); underline.fill.fore_color.rgb = SAGE_LIFT
    underline.line.fill.background(); underline.shadow.inherit = False

    # Three-beat value line
    add_closer(
        s,
        "Two minutes. Free. No signup.",
        Inches(7.3),
        accent_words={"Free.", "No", "signup."},
        size=24,
    )

    add_wordmark(s, center_y=Inches(9.9))


def slide_whats_next(prs):
    """Soft Phase 2 signal. Brief — doesn't pull attention from the ask."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  what’s next  · · ·", Inches(1.0))

    # Two-line headline
    tb, tf = add_text(s, Inches(0.6), Inches(2.6), CANVAS_W - Inches(1.2),
                      Inches(3.6), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.2
    r = p.add_run()
    set_run(r, text="Phase 1: ", font=SERIF, size=44, color=SLATE_LIFT)
    r = p.add_run()
    set_run(r, text="the audit.", font=SERIF, size=44, color=TEXT_LIGHT)

    p2 = add_paragraph(tf, alignment="center"); p2.line_spacing = 1.2
    r = p2.add_run()
    set_run(r, text="Phase 2: ", font=SERIF, size=44, color=SLATE_LIFT)
    r = p2.add_run()
    set_run(r, text="something bigger.", font=SERIF, size=44,
            color=PLUM_LIFT, italic=True)

    # Closer
    add_closer(
        s,
        "Stay close.",
        Inches(7.4),
        accent_words={"Stay", "close."},
        size=28,
    )

    add_wordmark(s, center_y=Inches(9.9))


# ─── Build ────────────────────────────────────────────────────────────

def build():
    # Re-use the deep deck's asset prerender (envelope, loupe, etc).
    prerender_sketches()

    prs = Presentation()
    prs.slide_width = CANVAS_W
    prs.slide_height = CANVAS_H

    slide_hook(prs)
    slide_why_now(prs)
    slide_proof(prs)
    slide_ask(prs)
    slide_whats_next(prs)

    out = "announcement.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    build()
