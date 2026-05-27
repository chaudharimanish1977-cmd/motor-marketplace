"""
Builds the 8-slide LinkedIn carousel for RightOffer's Beta launch — v2.

Differences vs v1:
  • Square 1080x1080 canvas — eliminates the bottom emptiness of 4:5
  • Dark mode palette (warm near-black bg + lifted plum/sage)
  • Brand sketch library woven in at slide-level rhythm
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import sketches as sk

# ─── Brand tokens (DARK MODE) ─────────────────────────────────────────
# Pitch deck dark bg from src/components/pitch-deck.tsx
BG_DARK    = RGBColor(0x13, 0x10, 0x0f)
# Brand v3 dark-mode flips
TEXT_LIGHT = RGBColor(0xf3, 0xee, 0xf0)   # was charcoal in light mode
PLUM_LIFT  = RGBColor(0xc4, 0x85, 0xc9)   # was #3a1e3d in light mode
SAGE_LIFT  = RGBColor(0xa8, 0xba, 0xa0)   # was #8b9d80
SLATE_LIFT = RGBColor(0x9a, 0x8f, 0x98)   # was #6b6571
HAIRLINE   = RGBColor(0x2a, 0x22, 0x28)   # subtle dark border

# SVG-friendly hex strings (resvg needs hex, not RGBColor)
HEX_TEXT  = "#f3eef0"
HEX_PLUM  = "#c485c9"
HEX_SAGE  = "#a8baa0"
HEX_SLATE = "#9a8f98"

SERIF = "Newsreader"
SERIF_FB = "Cambria"
MONO = "JetBrains Mono"
MONO_FB = "Consolas"

# Square: 10.8" × 10.8" → 1080×1080 at 100 ppi
CANVAS_W = Inches(10.8)
CANVAS_H = Inches(10.8)

ASSETS_DIR = os.path.join(os.path.dirname(__file__), "assets")
os.makedirs(ASSETS_DIR, exist_ok=True)


# ─── Sketch asset prerender ───────────────────────────────────────────
def _asset_path(name: str) -> str:
    return os.path.join(ASSETS_DIR, name + ".png")


def prerender_sketches() -> None:
    """Render each sketch we'll embed as a high-res PNG once, up front."""

    # Cars — render at ~4x the largest pptx use case for crispness.
    sk.render_to_png(sk.sketch_car_static(HEX_PLUM),
                     _asset_path("car-plum"), 1200, 600)
    sk.render_to_png(sk.sketch_car_static(HEX_TEXT),
                     _asset_path("car-text"), 600, 300)

    # Doc — light stroke on dark bg
    sk.render_to_png(sk.sketch_doc(HEX_TEXT),
                     _asset_path("doc"), 800, 624)
    sk.render_to_png(sk.sketch_doc(HEX_PLUM),
                     _asset_path("doc-plum"), 800, 624)

    # Loupe
    sk.render_to_png(sk.sketch_loupe(HEX_TEXT),
                     _asset_path("loupe"), 800, 624)
    sk.render_to_png(sk.sketch_loupe(HEX_PLUM),
                     _asset_path("loupe-plum"), 800, 624)

    # Verdict — plum lines, sage checks
    sk.render_to_png(sk.sketch_verdict(HEX_TEXT, HEX_SAGE),
                     _asset_path("verdict"), 800, 624)

    # Phase 2 trio — sage tone reads as "three competing insurers"
    sk.render_to_png(sk.sketch_sedan(HEX_SAGE),
                     _asset_path("sedan"), 1000, 625)
    sk.render_to_png(sk.sketch_hatchback(HEX_SAGE),
                     _asset_path("hatchback"), 1000, 625)
    sk.render_to_png(sk.sketch_suv(HEX_PLUM),
                     _asset_path("suv"), 1000, 625)

    # 5-point smileys. Render in light text colour; the active "5" reads
    # plum on its own slide to draw the eye.
    for r in (1, 2, 3, 4, 5):
        sk.render_to_png(sk.sketch_car_smiley(r, HEX_TEXT),
                         _asset_path(f"smiley-{r}"), 500, 500)
    # Plum-tone variants for the before/after slide
    sk.render_to_png(sk.sketch_car_smiley(2, HEX_TEXT),
                     _asset_path("smiley-2-light"), 600, 600)
    sk.render_to_png(sk.sketch_thank_you_car(HEX_PLUM),
                     _asset_path("thankyou"), 600, 720)

    # Motion strips — dashed road + speed lines
    sk.render_to_png(sk.sketch_motion_strip(HEX_SAGE, length=600),
                     _asset_path("motion-sage"), 1200, 100)
    sk.render_to_png(sk.sketch_motion_strip(HEX_PLUM, length=600),
                     _asset_path("motion-plum"), 1200, 100)

    # NEW — Envelope + mini-report + spotlight
    sk.render_to_png(sk.sketch_envelope(HEX_PLUM, HEX_SAGE),
                     _asset_path("envelope"), 800, 640)
    sk.render_to_png(sk.sketch_minireport(HEX_TEXT, HEX_SAGE),
                     _asset_path("minireport"), 720, 780)
    sk.render_to_png(sk.sketch_spotlight(HEX_SAGE),
                     _asset_path("spotlight-sage"), 600, 600)
    sk.render_to_png(sk.sketch_spotlight(HEX_PLUM),
                     _asset_path("spotlight-plum"), 600, 600)


# ─── Layout helpers ───────────────────────────────────────────────────

def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, CANVAS_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background(); bg.shadow.inherit = False
    return bg


def add_text(slide, left, top, width, height, *, align="center",
             anchor="middle"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[anchor]
    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    return tb, tf


def set_run(run, *, text, font, size, color, bold=False, italic=False,
            spacing=None):
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))


def add_paragraph(tf, *, alignment="center"):
    p = tf.add_paragraph()
    p.alignment = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[alignment]
    return p


def add_kicker(slide, text, top, *, color=None):
    """Top-line kicker. Bumped from 10pt to 14pt + bold so it actually
    reads at thumbnail size in feed."""
    if color is None:
        color = SAGE_LIFT
    tb, tf = add_text(slide, Inches(0.5), top, CANVAS_W - Inches(1.0),
                      Inches(0.55), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text=text.upper(), font=MONO, size=14, color=color,
            spacing=5, bold=True)
    return tb


def add_closer(slide, text, top, *, accent_words=None, color=None,
               size=20, italic=True, width=None):
    """Bottom-line body / closer. Bigger + brighter than the old 14-17pt
    slate-italic style. Optional `accent_words` list is rendered in
    sage to catch the eye.
    """
    if color is None:
        color = TEXT_LIGHT
    if width is None:
        width = CANVAS_W - Inches(1.4)
    left = (CANVAS_W - width) / 2
    tb, tf = add_text(slide, left, top, width, Inches(1.6),
                      align="center", anchor="top")
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.35
    # Strip punctuation when matching so callers can pass plain words.
    accents = {w.strip(".,;:!?").lower() for w in (accent_words or [])}
    import re
    tokens = re.findall(r"\S+\s*", text)
    for tok in tokens:
        word = tok.rstrip().strip(".,;:!?").lower()
        r = p.add_run()
        if word in accents:
            set_run(r, text=tok, font=SERIF, size=size,
                    color=SAGE_LIFT, italic=italic, bold=True)
        else:
            set_run(r, text=tok, font=SERIF, size=size,
                    color=color, italic=italic)
    return tb


def add_wordmark(slide, *, center_y):
    """italic plum 'r' + small-caps RIGHTOFFER + sage pill 'CAR'."""
    pill_w = Inches(0.55); pill_h = Inches(0.24)
    name_w = Inches(1.85)
    r_w = Inches(0.32)
    total_w = r_w + Inches(0.05) + name_w + Inches(0.15) + pill_w
    start_left = (CANVAS_W - total_w) / 2

    # italic plum r
    tb, tf = add_text(slide, start_left, center_y - Inches(0.22), r_w,
                      Inches(0.44), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="r", font=SERIF, size=26, color=PLUM_LIFT, italic=True)

    # RIGHTOFFER small-caps
    name_left = start_left + r_w + Inches(0.05)
    tb, tf = add_text(slide, name_left, center_y - Inches(0.18), name_w,
                      Inches(0.36), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="RIGHTOFFER", font=SERIF, size=14, color=TEXT_LIGHT,
            spacing=3)

    # sage pill
    pill_left = name_left + name_w + Inches(0.15)
    pill_top = center_y - pill_h / 2
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  pill_left, pill_top, pill_w, pill_h)
    pill.adjustments[0] = 0.5
    pill.fill.solid(); pill.fill.fore_color.rgb = SAGE_LIFT
    pill.line.fill.background(); pill.shadow.inherit = False
    tfp = pill.text_frame
    tfp.margin_left = tfp.margin_right = 0
    tfp.margin_top = tfp.margin_bottom = 0
    tfp.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = tfp.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    r = pp.add_run()
    set_run(r, text="CAR", font=MONO, size=9, color=BG_DARK, spacing=3,
            bold=True)


def add_pic(slide, asset, *, left, top, width):
    """Add a sketch image. Height is preserved by python-pptx if omitted."""
    return slide.shapes.add_picture(_asset_path(asset),
                                    left, top, width=width)


def add_centered_serif_block(slide, lines, *, top, height, size=48,
                             leading=1.1, color=None):
    if color is None:
        color = TEXT_LIGHT
    tb, tf = add_text(slide, Inches(0.6), top, CANVAS_W - Inches(1.2),
                      height, align="center", anchor="middle")
    for idx, (text, italic, override) in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        else:
            p = add_paragraph(tf, alignment="center")
        p.line_spacing = leading
        r = p.add_run()
        set_run(r, text=text, font=SERIF, size=size,
                color=override or color, italic=italic)


# ─── Slides ───────────────────────────────────────────────────────────

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)

    add_kicker(s, "· · ·  beta is live  · · ·", Inches(1.4))

    add_centered_serif_block(
        s,
        lines=[
            ("Your motor insurance,", False, TEXT_LIGHT),
            ("audited.", True, PLUM_LIFT),
        ],
        top=Inches(3.2), height=Inches(2.6), size=54, leading=1.1,
    )

    # Car sketch as cover-anchor — large plum-stroke + motion strip
    # underneath so the cover reads as "in motion" rather than parked.
    add_pic(s, "car-plum", left=Inches(3.4), top=Inches(6.3),
            width=Inches(4.0))
    add_pic(s, "motion-plum", left=Inches(2.4), top=Inches(7.7),
            width=Inches(6.0))

    # Subline — bumped from 18pt slate italic to 22pt text-light with
    # sage accents so it actually reads at thumbnail size.
    add_closer(
        s,
        "One email forward. Two minutes. Free.",
        Inches(8.6),
        accent_words={"Two", "minutes.", "Free."},
        size=22,
    )

    add_wordmark(s, center_y=Inches(9.9))


def slide_behavior(prs):
    """Phase-1 honest framing: we don't renew, we read.

    Previous copy ('We renew it. We don't read it.') implied we
    handled renewal, which Phase 1 does not. New framing puts the
    value squarely on us reading it for the customer.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  what we do  · · ·", Inches(0.9))

    # Two-line headline. Italic plum lands on "read" (slide 1) and the
    # whole subline (slide 2) for a strong rhythm.
    tb, tf = add_text(s, Inches(0.6), Inches(2.4),
                      CANVAS_W - Inches(1.2), Inches(2.8),
                      align="center", anchor="middle")
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER; p1.line_spacing = 1.1
    r = p1.add_run()
    set_run(r, text="We ", font=SERIF, size=58, color=TEXT_LIGHT)
    r = p1.add_run()
    set_run(r, text="read", font=SERIF, size=58, color=PLUM_LIFT, italic=True)
    r = p1.add_run()
    set_run(r, text=" it.", font=SERIF, size=58, color=TEXT_LIGHT)

    p2 = add_paragraph(tf, alignment="center"); p2.line_spacing = 1.1
    r = p2.add_run()
    set_run(r, text="So you don’t have to.", font=SERIF, size=44,
            color=TEXT_LIGHT, italic=True)

    # Doc sketch
    add_pic(s, "doc", left=Inches(4.05), top=Inches(5.6), width=Inches(2.7))

    # Closer — bigger, brighter, sage accent on "every line"
    add_closer(
        s,
        "Fifteen pages of fine print. We read every line and tell you what to ask.",
        Inches(8.5),
        accent_words={"every", "line."},
    )


def slide_stake(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  what’s at stake  · · ·", Inches(0.9))

    add_centered_serif_block(
        s,
        lines=[
            ("Most policies have", False, TEXT_LIGHT),
            ("a gap.", True, PLUM_LIFT),
        ],
        top=Inches(2.0), height=Inches(2.6), size=58, leading=1.1,
    )

    # Loupe — magnifying glass discovering gaps
    add_pic(s, "loupe", left=Inches(4.05), top=Inches(5.4),
            width=Inches(2.7))

    # Closer
    add_closer(
        s,
        "Zero-dep that expired. IDV set below market. Engine protector missing on a new car.",
        Inches(8.4),
        accent_words={"expired.", "below", "market."},
        size=19,
    )


def slide_wedge(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  the wedge  · · ·", Inches(0.9))

    # Headline — dropped 'column' (sounded editorial-newsroom);
    # 'honest read' lands better as a value claim.
    tb, tf = add_text(s, Inches(0.6), Inches(2.4),
                      CANVAS_W - Inches(1.2), Inches(2.2),
                      align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.1
    r = p.add_run()
    set_run(r, text="An ", font=SERIF, size=60, color=TEXT_LIGHT)
    r = p.add_run()
    set_run(r, text="honest", font=SERIF, size=60, color=PLUM_LIFT, italic=True)
    r = p.add_run()
    set_run(r, text=" read.", font=SERIF, size=60, color=TEXT_LIGHT)

    # Sub-line
    tb, tf = add_text(s, Inches(0.6), Inches(4.7), CANVAS_W - Inches(1.2),
                      Inches(0.7), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="Not a quote table.", font=SERIF, size=24,
            color=TEXT_LIGHT, italic=True)

    # Verdict sketch
    add_pic(s, "verdict", left=Inches(4.05), top=Inches(5.7),
            width=Inches(2.7))

    # Closer — bigger + sage accent on "no commission"
    add_closer(
        s,
        "No comparator. No commission. No insurer push. Just what your policy says.",
        Inches(8.6),
        accent_words={"No", "commission."},
    )


def slide_how(prs):
    """Horizontal three-step flow.

    Stations sit on a single row with dashed connectors between them
    so the eye reads forward → read → audit as a journey. Each station
    has a plum number badge above the sketch, a sage mono caption
    below.

    The MP4 animates the stations appearing in sequence.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  how it works  · · ·", Inches(0.9))

    tb, tf = add_text(s, Inches(0.5), Inches(1.9), CANVAS_W - Inches(1.0),
                      Inches(1.0), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="From inbox to ", font=SERIF, size=46, color=TEXT_LIGHT)
    r = p.add_run()
    set_run(r, text="audit", font=SERIF, size=46, color=PLUM_LIFT, italic=True)
    r = p.add_run()
    set_run(r, text=".", font=SERIF, size=46, color=TEXT_LIGHT)

    # Three stations, evenly spaced across canvas.
    stations = [
        ("doc",     "1", "FORWARD", "Email your policy"),
        ("loupe",   "2", "READ",    "Aryan reads it"),
        ("verdict", "3", "AUDIT",   "Lands in your inbox"),
    ]
    n = len(stations)
    station_w = Inches(2.6)
    gap = Inches(0.6)
    total_w = station_w * n + gap * (n - 1)
    start_x = (CANVAS_W - total_w) / 2

    station_top = Inches(4.4)
    sketch_w    = Inches(1.9)
    badge_d     = Inches(0.7)   # diameter of the number badge

    centers_y = station_top + Inches(1.9)  # vertical centerline of sketches

    # Dashed connectors between stations — implemented as a row of small
    # sage dot shapes (avoids the fragile connector + line.dash combo).
    n_dots = 6  # dots per gap
    dot_d = Inches(0.06)
    for i in range(n - 1):
        x_left  = start_x + (station_w + gap) * i + station_w
        x_right = x_left + gap
        gap_len = x_right - x_left
        spacing = gap_len / (n_dots + 1)
        for k in range(n_dots):
            cx = x_left + spacing * (k + 1)
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     cx - dot_d / 2,
                                     centers_y - dot_d / 2,
                                     dot_d, dot_d)
            dot.fill.solid(); dot.fill.fore_color.rgb = SAGE_LIFT
            dot.line.fill.background(); dot.shadow.inherit = False

    # Now the stations
    for i, (asset, num, label, caption) in enumerate(stations):
        x = start_x + (station_w + gap) * i

        # Plum number badge — sits above the sketch
        badge_left = x + (station_w - badge_d) / 2
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   badge_left, station_top,
                                   badge_d, badge_d)
        badge.fill.solid(); badge.fill.fore_color.rgb = PLUM_LIFT
        badge.line.fill.background(); badge.shadow.inherit = False
        tfb = badge.text_frame
        tfb.margin_left = tfb.margin_right = 0
        tfb.margin_top = tfb.margin_bottom = 0
        tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tfb.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run()
        set_run(r, text=num, font=SERIF, size=22, color=BG_DARK, bold=True)

        # Sketch
        sketch_left = x + (station_w - sketch_w) / 2
        add_pic(s, asset, left=sketch_left, top=station_top + Inches(0.95),
                width=sketch_w)

        # Sage uppercase label
        tb, tf = add_text(s, x, station_top + Inches(2.95), station_w,
                          Inches(0.4), align="center", anchor="middle")
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, text=label, font=MONO, size=13, color=SAGE_LIFT,
                spacing=5, bold=True)

        # Serif caption under
        tb, tf = add_text(s, x, station_top + Inches(3.45), station_w,
                          Inches(0.5), align="center", anchor="middle")
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, text=caption, font=SERIF, size=15, color=TEXT_LIGHT,
                italic=True)

    # Closer — sage accent words punch
    add_closer(
        s,
        "Two minutes from forward to audit. No app, no signup, no calls.",
        Inches(9.0),
        accent_words={"Two", "minutes"},
    )


def slide_get(prs):
    """Two-column. Left: numbered list of the five items.
    Right: an actual-looking 'RIGHTOFFER REPORT' card with the same
    items as labelled rows + sage ticks + a SCORE 62 badge.

    The intent is for the right side to feel like the deliverable, not
    a decoration — viewers should recognise it as 'oh, that's what I
    get.'
    """
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  what you’ll get  · · ·", Inches(0.9))

    # Headline
    tb, tf = add_text(s, Inches(0.5), Inches(1.9), CANVAS_W - Inches(1.0),
                      Inches(1.0), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="Five things, ", font=SERIF, size=42, color=TEXT_LIGHT)
    r = p.add_run()
    set_run(r, text="one report", font=SERIF, size=42, color=PLUM_LIFT,
            italic=True)
    r = p.add_run()
    set_run(r, text=".", font=SERIF, size=42, color=TEXT_LIGHT)

    # Two-column layout
    rows = [
        ("01", "Coverage snapshot"),
        ("02", "Gap analysis"),
        ("03", "IDV cross-check"),
        ("04", "Add-on relevance"),
        ("05", "Renewal calendar"),
    ]
    list_left = Inches(0.9)
    list_top  = Inches(3.7)
    row_h     = Inches(0.85)
    list_w    = Inches(4.6)

    for i, (num, text) in enumerate(rows):
        y = list_top + Emu(int(row_h) * i)
        # Numeral in mono sage
        tb, tf = add_text(s, list_left, y, Inches(0.8), Inches(0.8),
                          align="left", anchor="middle")
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        set_run(r, text=num, font=MONO, size=16, color=SAGE_LIFT,
                bold=True, spacing=3)
        # Label
        tb, tf = add_text(s, list_left + Inches(0.85), y, list_w, Inches(0.8),
                          align="left", anchor="middle")
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        set_run(r, text=text, font=SERIF, size=22, color=TEXT_LIGHT)

    # Mini-report card on the right
    add_pic(s, "minireport", left=Inches(6.1), top=Inches(3.4),
            width=Inches(4.0))

    # Closer
    add_closer(
        s,
        "One audit. One report. Everything you need to renew with confidence.",
        Inches(9.2),
        accent_words={"confidence."},
        size=18,
    )


def slide_scale(prs):
    """Sits between 'Five things' and 'Phase 2'.
    Shows the 5-point CarSmiley scale the audit uses to grade each gap.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  how we rate  · · ·", Inches(0.9))

    # Headline
    add_centered_serif_block(
        s,
        lines=[
            ("Every finding,", False, TEXT_LIGHT),
            ("on a 5-point scale.", True, PLUM_LIFT),
        ],
        top=Inches(1.7), height=Inches(2.4), size=44, leading=1.15,
    )

    # Row of 5 smileys, each on a sage spotlight glow. The active one
    # ('Delighted' = 5) gets a plum spotlight to draw the eye.
    smiley_top = Inches(4.4)
    smiley_w   = Inches(1.5)
    spot_w     = Inches(1.9)   # spotlight is wider than the smiley
    gap        = Inches(0.45)
    total      = smiley_w * 5 + gap * 4
    start_left = (CANVAS_W - total) / 2
    spot_offset = (smiley_w - spot_w) / 2  # negative — spotlight sits behind

    labels = [
        ("1", "Awful"),
        ("2", "Bad"),
        ("3", "Okay"),
        ("4", "Good"),
        ("5", "Delighted"),
    ]
    for i, (num, label) in enumerate(labels):
        x = start_left + (smiley_w + gap) * i
        spotlight_asset = "spotlight-plum" if (i == 4) else "spotlight-sage"
        # Spotlight first so smiley sits on top
        add_pic(s, spotlight_asset,
                left=x + spot_offset, top=smiley_top + Inches(-0.2),
                width=spot_w)
        add_pic(s, f"smiley-{i + 1}", left=x, top=smiley_top, width=smiley_w)

        # Mono numeral
        tb, tf = add_text(s, x, smiley_top + Inches(1.65), smiley_w,
                          Inches(0.4), align="center", anchor="middle")
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, text=num, font=MONO, size=16, color=SAGE_LIFT,
                spacing=3, bold=True)

        # Label
        is_active = (i == 4)
        tb, tf = add_text(s, x, smiley_top + Inches(2.05), smiley_w,
                          Inches(0.4), align="center", anchor="middle")
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, text=label, font=SERIF, size=16,
                color=PLUM_LIFT if is_active else TEXT_LIGHT,
                italic=True, bold=is_active)

    # Closer
    add_closer(
        s,
        "Every gap gets a face. You see exactly where you stand.",
        Inches(8.6),
        accent_words={"face.", "stand."},
        size=20,
    )


def slide_before_after(prs):
    """Editorial before / after.

    'From confused' lands first (read order), then a strongly
    emphasised 'to CLEAR.' with a sage underline so the eye lands hard
    on the destination state. In the MP4 these stagger in.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  before & after  · · ·", Inches(0.9))

    # Headline line 1 — "From confused"
    tb, tf = add_text(s, Inches(0.6), Inches(2.0), CANVAS_W - Inches(1.2),
                      Inches(1.1), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="From confused", font=SERIF, size=48, color=TEXT_LIGHT)

    # Headline line 2 — "to CLEAR." with extra-large emphasis on CLEAR
    tb, tf = add_text(s, Inches(0.6), Inches(3.0), CANVAS_W - Inches(1.2),
                      Inches(1.6), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="to ", font=SERIF, size=48, color=TEXT_LIGHT)
    r = p.add_run()
    set_run(r, text="clear", font=SERIF, size=72, color=PLUM_LIFT,
            italic=True, bold=True)
    r = p.add_run()
    set_run(r, text=".", font=SERIF, size=72, color=PLUM_LIFT)

    # Sage underline beneath 'clear' for spotlight effect
    underline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(4.45), Inches(4.55),
                                   Inches(2.4), Emu(36000))  # ~3pt tall
    underline.fill.solid(); underline.fill.fore_color.rgb = SAGE_LIFT
    underline.line.fill.background(); underline.shadow.inherit = False

    # Two halves — sad smiley left, thank-you car right
    half_top = Inches(5.2)
    img_w    = Inches(2.6)
    left_x   = Inches(1.7)
    right_x  = CANVAS_W - left_x - img_w
    smiley_w = Inches(2.6)

    add_pic(s, "smiley-2-light", left=left_x, top=half_top + Inches(0.3),
            width=smiley_w)
    add_pic(s, "thankyou", left=right_x, top=half_top, width=img_w)

    # Captions under each
    cap_top = Inches(8.2)

    # LEFT — Before
    tb, tf = add_text(s, left_x - Inches(0.4), cap_top, img_w + Inches(0.8),
                      Inches(0.4), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="BEFORE", font=MONO, size=11, color=SAGE_LIFT, spacing=4)

    tb, tf = add_text(s, left_x - Inches(0.4), cap_top + Inches(0.45),
                      img_w + Inches(0.8), Inches(0.6),
                      align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="Fifteen pages of fine print.",
            font=SERIF, size=15, color=SLATE_LIFT, italic=True)

    # RIGHT — After
    tb, tf = add_text(s, right_x - Inches(0.4), cap_top, img_w + Inches(0.8),
                      Inches(0.4), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="AFTER", font=MONO, size=11, color=SAGE_LIFT, spacing=4)

    tb, tf = add_text(s, right_x - Inches(0.4), cap_top + Inches(0.45),
                      img_w + Inches(0.8), Inches(0.6),
                      align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="A two-minute audit.",
            font=SERIF, size=15, color=TEXT_LIGHT, italic=True)


def slide_phase2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  what’s brewing  · · ·", Inches(0.9))

    add_centered_serif_block(
        s,
        lines=[
            ("Phase 2 is", False, TEXT_LIGHT),
            ("brewing.", True, PLUM_LIFT),
        ],
        top=Inches(1.8), height=Inches(2.6), size=56, leading=1.1,
    )

    # Trio of cars — implied "insurers competing"
    car_top = Inches(5.0)
    add_pic(s, "sedan",     left=Inches(0.4),  top=car_top, width=Inches(3.2))
    add_pic(s, "hatchback", left=Inches(3.8),  top=car_top, width=Inches(3.2))
    add_pic(s, "suv",       left=Inches(7.2),  top=car_top, width=Inches(3.2))
    # Motion strip running under all three — reads as "in motion / racing"
    add_pic(s, "motion-sage", left=Inches(0.4), top=Inches(7.4),
            width=Inches(10.0))

    # Closer — fun, intriguing, no Phase-2 reveals
    add_closer(
        s,
        "We’re building something that’ll change how you renew. Stay close.",
        Inches(8.6),
        accent_words={"Stay", "close."},
        size=20,
    )


def slide_cta(prs):
    """Email-primary CTA.

    The email address is the hero — it's the action we want. The
    envelope sketch sits above it, the website URL drops to a small
    secondary line below.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_kicker(s, "· · ·  forward your policy  · · ·", Inches(0.9))

    # Envelope sketch — the hero visual
    add_pic(s, "envelope", left=Inches(4.0), top=Inches(2.0),
            width=Inches(2.8))

    # "Forward to" tag — small mono sage above the email
    tb, tf = add_text(s, Inches(0.5), Inches(4.8), CANVAS_W - Inches(1.0),
                      Inches(0.5), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="FORWARD TO", font=MONO, size=13, color=SAGE_LIFT,
            spacing=5, bold=True)

    # Email — primary CTA, BIG in plum-on-cream feel
    tb, tf = add_text(s, Inches(0.4), Inches(5.4), CANVAS_W - Inches(0.8),
                      Inches(1.5), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="review@rightoffer.in", font=SERIF, size=46,
            color=PLUM_LIFT, italic=True, bold=True)

    # Sage underline under the email — spotlight effect
    underline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1.5), Inches(6.55),
                                   Inches(7.8), Emu(36000))
    underline.fill.solid(); underline.fill.fore_color.rgb = SAGE_LIFT
    underline.line.fill.background(); underline.shadow.inherit = False

    # Closer — value reminder, bigger than before
    add_closer(
        s,
        "Two minutes. No signup. The audit is on us.",
        Inches(7.2),
        accent_words={"audit", "us."},
        size=22,
    )

    # Tiny website URL as secondary
    tb, tf = add_text(s, Inches(0.5), Inches(8.8), CANVAS_W - Inches(1.0),
                      Inches(0.5), align="center", anchor="middle")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text="rightoffer.in", font=MONO, size=13,
            color=SLATE_LIFT, spacing=3)

    add_wordmark(s, center_y=Inches(9.9))


# ─── Build ────────────────────────────────────────────────────────────

def build():
    prerender_sketches()

    prs = Presentation()
    prs.slide_width = CANVAS_W
    prs.slide_height = CANVAS_H

    slide_cover(prs)
    slide_behavior(prs)
    slide_stake(prs)
    slide_wedge(prs)
    slide_how(prs)
    slide_get(prs)
    slide_scale(prs)         # NEW — 5-point CarSmiley row
    slide_before_after(prs)  # NEW — confused → clear (happiness)
    slide_phase2(prs)
    slide_cta(prs)

    out = "beta-launch-carousel.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    build()
